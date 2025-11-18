import streamlit as st
import pdfplumber
import pptx
import google.generativeai as genai
import json
import re
import urllib.parse

# --- CONFIGURATION ---
st.set_page_config(page_title="Eureka Pitch Scorer", layout="wide")

# --- 1. AUTHENTICATION ---
api_key = st.secrets.get("GEMINI_API_KEY")
if not api_key:
    st.error("🔑 API Key missing. Please add GEMINI_API_KEY to Streamlit Secrets.")
    st.stop()

genai.configure(api_key=api_key)

# --- 2. KNOWLEDGE BASE (THE ANCHORS) ---
RUBRIC_GUIDE = """
CASE STUDY ANCHORS (Use these to grade):

1. PROBLEM IDENTIFICATION ("Why you?")
   ⭐ 1 STAR (BAD): "We are passionate students who love music." 
   (Reasoning: Generic passion, no unique leverage.)
   ⭐ 3 STAR (GOOD): "Our CTO holds a patent in audio signal processing and I managed a $2M inventory at Guitar Center."
   (Reasoning: Specific, verifiable, relevant domain expertise.)

2. CUSTOMER DISCOVERY ("Who is the customer?")
   ⭐ 1 STAR (BAD): "Everyone who owns a home is our customer."
   (Reasoning: TAM is not a customer profile. Too broad.)
   ⭐ 3 STAR (GOOD): "Our beachhead is single-family homeowners in the Northeast with oil heat (12% of region)."
   (Reasoning: Specific geography, demographic, and technical constraint.)

3. VALIDATION ("How do you know?")
   ⭐ 1 STAR (BAD): "We sent out a survey and people liked it."
   (Reasoning: Surveys are weak evidence of purchase intent.)
   ⭐ 3 STAR (GOOD): "We pre-sold 50 units at $20 each using a smoke-test landing page."
   (Reasoning: Financial commitment and actual behavior tracked.)
"""

RUBRIC_QUESTIONS = """
SECTION 1: PROBLEM
1. What is the problem you want to solve?
2. Why do you care about solving this problem?
3. Why are you uniquely qualified to solve this problem?
4. Initial thoughts on finding customers?
5. Initial thoughts on monetization?

SECTION 2: DISCOVERY
6. Who is the customer/end user?
7. Have you carved out user profile specifics? (Who have you talked to?)
8. How are customers currently solving this problem?
9. What are the competitive products?
10. Prototype status?
11. Analysis of results?
12. What do you need now?
"""

# --- 3. HELPER FUNCTIONS ---
def extract_text(file, file_type):
    text = ""
    try:
        if file_type == "pdf":
            with pdfplumber.open(file) as pdf:
                text = "\n".join([p.extract_text() or "" for p in pdf.pages])
        elif file_type == "pptx":
            prs = pptx.Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception:
        return None
    return text

def clean_json_response(response_text):
    text = re.sub(r'```json\n?', '', response_text)
    text = re.sub(r'```', '', text)
    return text.strip()

def generate_google_link(query):
    encoded = urllib.parse.quote(query)
    return f"https://www.google.com/search?q={encoded}"

# --- 4. AGENT 1: THE JUDGE ---
def analyze_pitch(deck_text):
    prompt = f"""
    You are a strict Judge for the 'Eureka' Pitch Competition.
    
    TASK: Score the uploaded pitch deck based on the RUBRIC below.
    Use the "CASE STUDY ANCHORS" to determine if a score is 1 or 3.

    INPUT PITCH DECK:
    "{deck_text[:30000]}"

    RUBRIC QUESTIONS:
    {RUBRIC_QUESTIONS}

    CASE STUDY ANCHORS (STRICT RULES):
    {RUBRIC_GUIDE}

    OUTPUT FORMAT:
    Respond with VALID JSON ONLY:
    {{
        "reviews": [
            {{
                "question": "1. What is the problem?",
                "score": 1,
                "reasoning": "The deck only states..."
            }},
            ...
        ],
        "total_score": 0,
        "hard_truth": "Summary paragraph."
    }}
    """
    
    # Try models in order (Using your approved list)
    model_options = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-1.5-pro", "gemini-1.5-flash"]
    
    for m in model_options:
        try:
            # Try strict JSON mode
            model = genai.GenerativeModel(m, generation_config={"response_mime_type": "application/json"})
            return model.generate_content(prompt).text
        except:
            try:
                # Retry with 'models/' prefix if needed
                model = genai.GenerativeModel(f"models/{m}", generation_config={"response_mime_type": "application/json"})
                return model.generate_content(prompt).text
            except:
                continue
    return None

# --- 5. AGENT 2: THE TEACHER (FIXED MODEL SELECTION) ---
def get_case_studies(weak_areas_list):
    # We turn the list of weak questions into a string
    weaknesses_str = "\n".join([f"- {w['question']} (Score: {w['score']})" for w in weak_areas_list])
    
    prompt = f"""
    You are a Startup Mentor. The user has failed the following areas in their pitch:
    {weaknesses_str}
    
    TASK:
    For EACH weakness, identify a famous successful startup (Airbnb, Dropbox, Uber, DoorDash, etc.) that solved this specific problem perfectly in their early pitch deck.
    
    OUTPUT FORMAT (JSON):
    {{
        "case_studies": [
            {{
                "weakness": "Customer Discovery",
                "example_company": "Airbnb",
                "lesson": "Airbnb didn't just say 'travelers'. They specifically targeted attendees of a design conference in SF when hotels were sold out.",
                "search_query": "Airbnb pitch deck customer validation slide" 
            }}
        ]
    }}
    """
    
    # FIXED: We now loop through the approved models instead of hardcoding 1.5-flash
    model_options = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-1.5-pro", "gemini-1.5-flash"]
    
    for m in model_options:
        try:
            model = genai.GenerativeModel(m, generation_config={"response_mime_type": "application/json"})
            return model.generate_content(prompt).text
        except:
             try:
                # Retry with 'models/' prefix
                model = genai.GenerativeModel(f"models/{m}", generation_config={"response_mime_type": "application/json"})
                return model.generate_content(prompt).text
             except:
                continue
    return None

# --- 6. THE UI ---
st.title("💡 Eureka Pitch Scorer & Coach")

if "analysis_data" not in st.session_state:
    st.session_state["analysis_data"] = None

uploaded_file = st.file_uploader("Upload Pitch Deck", type=["pdf", "pptx"])

if uploaded_file and st.button("Run Evaluation"):
    with st.spinner("Reading file..."):
        ftype = uploaded_file.name.split(".")[-1].lower()
        extracted_text = extract_text(uploaded_file, ftype)
        
    if extracted_text:
        with st.spinner("Judging against Anchors..."):
            raw_result = analyze_pitch(extracted_text)
            if raw_result:
                try:
                    st.session_state["analysis_data"] = json.loads(clean_json_response(raw_result))
                except:
                    st.error("Error parsing AI response. Please try again.")

# --- DISPLAY RESULTS ---
if st.session_state["analysis_data"]:
    data = st.session_state["analysis_data"]
    
    # Top Section: Score
    col1, col2 = st.columns([1, 3])
    with col1:
        st.metric("Total Score", f"{data.get('total_score')}/36")
    with col2:
        st.error(f"**The Hard Truth:** {data.get('hard_truth')}")

    st.divider()
    st.subheader("📋 Detailed Report Card")
    
    # Report Card Loop
    for review in data['reviews']:
        score = review['score']
        color = "red" if score == 1 else "orange" if score == 2 else "green"
        
        with st.container():
            c1, c2 = st.columns
